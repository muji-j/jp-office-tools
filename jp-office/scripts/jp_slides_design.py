# -*- coding: utf-8 -*-
"""jp-slides 描画ツールキット。

SP3 で承認された16テーマ・モダンデザインの共通ヘルパー群(矩形・楕円・多角形・波形・
ピル・カード各種・テキスト・罫線)。承認済みプロトタイプ(design_ref/sp3_rich.py)から
座標・色・構成を変えずに移植したもので、以降のタスクで render_cover 等から利用する。

このモジュールはレンダリング時にのみインポートされる想定のため、python-pptx は
モジュール冒頭で直接インポートする(未インストール環境向けの遅延インポートは
呼び出し側の jp_slides.py の _require_pptx が担う)。
"""
from __future__ import annotations

import math

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

SW, SH = 13.333, 7.5
RR = MSO_SHAPE.ROUNDED_RECTANGLE
RTOP = MSO_SHAPE.ROUND_2_SAME_RECTANGLE
GO = "游ゴシック"


def rgb(h):
    return RGBColor.from_string(h)


def new_prs():
    p = Presentation()
    p.slide_width = Inches(SW)
    p.slide_height = Inches(SH)
    return p


def slide(prs, bg):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = rgb(bg)
    return s


def _ea(run, name):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def _spc(run, val):
    try:
        run._r.get_or_add_rPr().set("spc", str(int(val)))
    except Exception:
        pass


def _alpha(sp, pct):
    try:
        el = sp.fill.fore_color._xFill.find(qn('a:srgbClr'))
        a = el.makeelement(qn('a:alpha'), {'val': str(int(pct * 1000))})
        el.append(a)
    except Exception:
        pass


def rect(s, x, y, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE,
          alpha=None, radius=None, rot=None, lalpha=None):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if rot:
        sp.rotation = rot
    if radius is not None:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = rgb(fill)
        if alpha is not None:
            _alpha(sp, alpha)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = rgb(line)
        sp.line.width = Pt(lw)
        if lalpha is not None:
            try:
                el = sp.line.color._xFill.find(qn('a:srgbClr'))
                a = el.makeelement(qn('a:alpha'), {'val': str(int(lalpha * 1000))})
                el.append(a)
            except Exception:
                pass
    return sp


def oval(s, x, y, w, h, fill=None, line=None, lw=1.0, alpha=None):
    return rect(s, x, y, w, h, fill=fill, line=line, lw=lw, shape=MSO_SHAPE.OVAL, alpha=alpha)


def poly(s, pts, fill=None, alpha=None, line=None, lw=1.0):
    P = [(int(x * 914400), int(y * 914400)) for x, y in pts]
    fb = s.shapes.build_freeform(P[0][0], P[0][1], scale=1.0)
    fb.add_line_segments(P[1:], close=True)
    sp = fb.convert_to_shape()
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = rgb(fill)
        if alpha is not None:
            _alpha(sp, alpha)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = rgb(line)
        sp.line.width = Pt(lw)
    return sp


def wave(s, x0, x1, ybase, amp, periods, fill, alpha=None, down=True, samples=64):
    pts = []
    for i in range(samples + 1):
        t = i / samples
        x = x0 + (x1 - x0) * t
        y = ybase + amp * math.sin(2 * math.pi * periods * t + math.pi)
        pts.append((x, y))
    pts += [(x1, SH + 0.6 if down else -0.6), (x0, SH + 0.6 if down else -0.6)]
    return poly(s, pts, fill=fill, alpha=alpha)


def R(txt, font, size, color, bold=False, spc=0):
    return (txt, font, size, color, bold, spc)


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, ls=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    first = True
    for line in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if ls:
            p.line_spacing = ls
        for (t, f, sz, c, b, sp) in line:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(sz)
            r.font.bold = b
            r.font.color.rgb = rgb(c)
            _ea(r, f)
            if sp:
                _spc(r, sp)
    return tb


def pill(s, x, y, w, h, txt, fill, fg, size=11, line=None, lw=1.0):
    rect(s, x, y, w, h, fill=fill, line=line, lw=lw, shape=RR, radius=0.5)
    text(s, x, y, w, h, [[R(txt, GO, size, fg, True, 60)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def hline(s, x, y, w, color, pt=1.0, alpha=None):
    rect(s, x, y, w, pt / 72.0, fill=color, alpha=alpha)


def vline(s, x, y, h, color, pt=1.0, alpha=None):
    rect(s, x, y, pt / 72.0, h, fill=color, alpha=alpha)


# ---- カードバリエーション ----
def card_plain(s, x, y, w, h, fill, shadow=None, radius=0.06, line=None, lw=0.75, alpha=None, lalpha=None):
    if shadow:
        shp = rect(s, x + 0.05, y + 0.09, w, h, fill=shadow, shape=RR, radius=radius)
        _alpha(shp, 50)
    return rect(s, x, y, w, h, fill=fill, shape=RR, radius=radius, line=line, lw=lw, alpha=alpha, lalpha=lalpha)


def card_topstrip(s, x, y, w, h, body, strip, shadow=None, radius=0.06, sh_h=0.16):
    card_plain(s, x, y, w, h, body, shadow=shadow, radius=radius)
    rect(s, x, y, w, sh_h, fill=strip, shape=RTOP, radius=0.5)


def card_leftbar(s, x, y, w, h, body, bar, shadow=None, radius=0.06):
    card_plain(s, x, y, w, h, body, shadow=shadow, radius=radius)
    rect(s, x + 0.28, y + 0.3, 0.1, h - 0.6, fill=bar, shape=RR, radius=0.5)


def card_outline(s, x, y, w, h, line, lw=1.25, radius=0.06, fill=None, alpha=None):
    return rect(s, x, y, w, h, fill=fill, alpha=alpha, line=line, lw=lw, shape=RR, radius=radius)


def card_split(s, x, y, w, h, body, divider, shadow=None, radius=0.06):
    card_plain(s, x, y, w, h, body, shadow=shadow, radius=radius)
    hline(s, x + 0.4, y + h * 0.5, w - 0.8, divider, 1.0)


# ---- 色ユーティリティ ----
def _tint(hexcolor, amount=0.85):
    """アクセント色を白側にブレンドした淡色16進を返す(ゴースト数字などの背景装飾用)。

    design_ref/sp3_themes.py はゴースト数字などの淡色をテーマごとに直接ハード
    コードしていたが、ここではプロファイルのアクセント色から機械的に導出する。
    """
    r = int(hexcolor[0:2], 16)
    g = int(hexcolor[2:4], 16)
    b = int(hexcolor[4:6], 16)
    r = int(r + (255 - r) * amount)
    g = int(g + (255 - g) * amount)
    b = int(b + (255 - b) * amount)
    return f"{r:02X}{g:02X}{b:02X}"


C = PP_ALIGN.CENTER
L = PP_ALIGN.LEFT
MID = MSO_ANCHOR.MIDDLE


# ================= シグネチャ背景 16種 =================
# variant: "cover"(フル強度) / "body"(控えめ) / "section"(フル、またはテーマ固有の変化形)。
# 各関数は design_ref/sp3_themes.py の対応テーマ関数から背景・装飾図形の部分だけを
# 抽出して移植し、ハードコードされた色をプロファイルのパレットキーに置き換えたもの。

def _bg_airyrule(slide, prof, variant):
    """霞(kasumi) — 縦の細罫線 + セクション見出し下の短いアクセント罫線。"""
    if variant == "cover":
        vline(slide, 4.7, 1.1, 5.3, prof["rule"], 1.0)
    elif variant == "section":
        hline(slide, SW / 2 - 1.25, 4.55, 2.5, prof["accent"], 1.4)


def _bg_mono(slide, prof, variant):
    """白磁(hakuji) — 装飾図形なし。余白とタイポグラフィのみで見せるテーマ。"""
    return


def _bg_rail(slide, prof, variant):
    """石板(sekiban) — 左側の丸みレール。section は幅広版のレール。"""
    if variant == "section":
        rect(slide, -0.4, 0.35, 5.6, 6.8, fill=prof["accent"], shape=RR, radius=0.05)
    else:
        rect(slide, -0.4, 0.35, 2.1, 6.8, fill=prof["accent"], shape=RR, radius=0.06)


def _bg_panel(slide, prof, variant):
    """藍(ai) — section はアクセント全面。cover/body は背景装飾なし。

    原典 design_ref/sp3_themes.py の ai() は cover/body に背景装飾を持たず、
    リッチさはコンテンツカード側で担保している。ここでもプロファイルの背景色
    (slide() で塗った bg)のみを残し、section の全面アクセントだけ原典どおり描く。
    """
    if variant == "section":
        rect(slide, 0, 0, SW, SH, fill=prof["accent"])


def _bg_diagonal(slide, prof, variant):
    """常磐(tokiwa) — 対角フリーフォームのウェッジ。variant ごとに大きさが変わる。"""
    accent = prof["accent"]
    if variant == "cover":
        poly(slide, [(9.4, 0), (SW, 0), (SW, SH), (8.2, SH)], fill=accent)
        poly(slide, [(9.05, 0), (9.3, 0), (7.95, SH), (7.7, SH)], fill=accent, alpha=26)
    elif variant == "body":
        poly(slide, [(12.6, 0), (SW, 0), (SW, 0.85)], fill=accent)
    elif variant == "section":
        poly(slide, [(5.6, 0), (SW, 0), (SW, SH), (4.5, SH)], fill=accent)


def _bg_band(slide, prof, variant):
    """鉄紺(tetsukon) — 上部ヘッダーバンド。section は全面バンド。"""
    navy = prof["accent"]
    line = prof.get("accent2") or prof["accent"]
    if variant == "cover":
        rect(slide, 0, 0, SW, 2.5, fill=navy)
        rect(slide, 0, 2.5, SW, 0.06, fill=line)
    elif variant == "body":
        rect(slide, 0, 0, SW, 0.9, fill=navy)
    elif variant == "section":
        rect(slide, 0, 0, SW, SH, fill=navy)


def _bg_lines(slide, prof, variant):
    """青碧(seiheki) — 罫線による画面の比率分割 + 右側カラーフィールド。"""
    accent = prof["accent"]
    rule = prof["rule"]
    if variant == "cover":
        rect(slide, 8.25, 0, SW - 8.25, SH, fill=accent)
        vline(slide, 8.25, 0, SH, prof.get("accent2") or prof["ink"], 2.0)
        hline(slide, 0.9, 4.35, 7.0, rule, 1.2)
    elif variant == "body":
        hline(slide, 0.9, 2.25, 11.5, rule, 1.4)
    elif variant == "section":
        rect(slide, 0, 0, 5.2, SH, fill=accent)


def _bg_circle(slide, prof, variant):
    """朱(shu) — 大型クロップ円 + 重なる小円。section はアクセント全面。"""
    accent = prof["accent"]
    accent2 = prof.get("accent2") or accent
    if variant == "cover":
        oval(slide, 8.7, 3.3, 6.4, 6.4, fill=accent)
        oval(slide, 7.7, 4.6, 1.5, 1.5, fill=accent2, alpha=80)
    elif variant == "section":
        rect(slide, 0, 0, SW, SH, fill=accent)


def _bg_ghostnum(slide, prof, variant):
    """山吹(yamabuki) — 淡色オーバーサイズのゴースト数字。

    design_ref では x+w がスライド幅を超えて右にブリードしていたが、ここでは
    画面内に収まるよう座標を調整した修正版。
    """
    ghost = _tint(prof["accent"], 0.85)
    if variant == "cover":
        text(slide, 8.9, 1.35, 4.2, 2.6, [[R("8", GO, 118, ghost, True)]])
    elif variant == "section":
        text(slide, 7.9, 1.4, 5.0, 3.2, [[R("8", GO, 168, ghost, True)]])


def _bg_colorblock(slide, prof, variant):
    """彩層(saisou) — バウハウス風の色ブロック。section は上部ブロック+区切り線。"""
    ind = prof["accent"]
    gold = prof.get("accent2") or prof["muted"]
    if variant == "cover":
        rect(slide, 8.6, 0, SW - 8.6, SH, fill=ind)
        rect(slide, 8.6, 4.6, SW - 8.6, 2.9, fill=gold)
        oval(slide, 10.0, 1.0, 2.2, 2.2, fill=gold)
        poly(slide, [(8.6, 3.2), (10.4, 3.2), (8.6, 5.0)], fill="FFFFFF")
    elif variant == "section":
        rect(slide, 0, 0, SW, 3.4, fill=ind)
        rect(slide, 0, 3.4, SW, 0.1, fill=gold)


def _bg_glassdark(slide, prof, variant):
    """墨(sumi) — ダーク背景のみ。cover/body に背景装飾はない。

    原典 design_ref/sp3_themes.py の sumi() はヘアラインもグロー円も持たず、
    ネオン感はカード枠線・アクセント文字色だけで表現している。section のみ、
    原典どおり見出し下の短いアクセントバーを描く。
    """
    if variant == "section":
        rect(slide, 0.95, 4.6, 0.9, 0.08, fill=prof["accent"])


def _bg_horizon(slide, prof, variant):
    """藍鉄(aitetsu) — 大型楕円のホライズン + グロー円。全 variant で共通。"""
    base = prof.get("card") or prof["bg"]
    accent = prof["accent"]
    oval(slide, -6, -9.5, 26, 14, fill=base)
    oval(slide, 8.0, -2.4, 7.5, 7.5, fill=accent, alpha=12)
    oval(slide, 10.2, -0.8, 3.6, 3.6, fill=accent, alpha=17)


def _bg_wave(slide, prof, variant):
    """桜(sakura) — サイン波。section は上向きの波(down=False)。"""
    rose = prof["accent"]
    if variant == "cover":
        wave(slide, 0, SW, 6.0, 0.55, 1.2, rose, alpha=30)
        wave(slide, 0, SW, 6.55, 0.4, 1.6, rose, alpha=55)
        oval(slide, 10.6, 0.7, 3.3, 3.3, fill=rose, alpha=22)
    elif variant == "body":
        wave(slide, 0, SW, 6.7, 0.35, 1.4, rose, alpha=22)
    elif variant == "section":
        wave(slide, 0, SW, 3.3, 0.7, 1.1, rose, alpha=32, down=False)


def _bg_bottomband(slide, prof, variant):
    """亜麻(ama) — 下部の波/バンド。"""
    terra = prof["accent"]
    if variant == "cover":
        wave(slide, 0, SW, 5.3, 0.45, 1.1, terra, alpha=88)
        wave(slide, 0, SW, 5.7, 0.35, 1.5, terra, alpha=55)
    elif variant == "body":
        rect(slide, 0, 6.85, SW, 0.65, fill=terra)
    elif variant == "section":
        rect(slide, 0, 4.2, SW, 3.3, fill=terra)


def _bg_frame(slide, prof, variant):
    """藤(fuji) — 全周ヘアラインフレーム。section は二重フレーム。"""
    wis = prof["accent"]
    rect(slide, 0.35, 0.35, SW - 0.7, SH - 0.7, line=wis, lw=0.75, lalpha=60)
    if variant == "section":
        rect(slide, 0.6, 0.6, SW - 1.2, SH - 1.2, line=wis, lw=0.5, lalpha=40)


def _bg_spine(slide, prof, variant):
    """明朝(mincho) — 縦のスパイン罫線。variant ごとに位置が変わる。"""
    rule = prof["rule"]
    if variant == "cover":
        vline(slide, 4.6, 0, SH, rule, 0.75)
    elif variant == "body":
        vline(slide, 0.9, 1.6, 4.8, rule, 0.75)
    elif variant == "section":
        vline(slide, 6.66, 0, SH, rule, 0.75)


_BG_DISPATCH = {
    "airyrule": _bg_airyrule, "mono": _bg_mono, "rail": _bg_rail, "panel": _bg_panel,
    "diagonal": _bg_diagonal, "band": _bg_band, "lines": _bg_lines, "circle": _bg_circle,
    "ghostnum": _bg_ghostnum, "colorblock": _bg_colorblock, "glassdark": _bg_glassdark,
    "horizon": _bg_horizon, "wave": _bg_wave, "bottomband": _bg_bottomband,
    "frame": _bg_frame, "spine": _bg_spine,
}


def _bg(slide, prof, variant):
    """プロファイルの bgsig に応じてシグネチャ背景を描画するディスパッチャ。

    variant は "cover" / "body" / "section" のいずれか。
    """
    fn = _BG_DISPATCH.get(prof.get("bgsig"))
    if fn is None:
        raise ValueError(f"未知の背景シグネチャです: {prof.get('bgsig')!r}")
    fn(slide, prof, variant)


# ---- キッカー(見出し上のラベル装飾) ----
def _kicker(slide, prof, x, y, w, h, label):
    if not label:
        return
    style = prof.get("kicker", "plain")
    accent = prof["accent"]
    if style == "pill":
        soft = prof.get("card") if prof.get("dark") else _tint(accent, 0.85)
        pill(slide, x, y, w, h, label, soft, accent, 12)
    elif style == "pill_outline":
        rect(slide, x, y, w, h, line=accent, lw=1.0, shape=RR, radius=0.5)
        text(slide, x, y, w, h, [[R(label, GO, 11, accent, True, 80)]], align=C, anchor=MID)
    elif style == "underline":
        text(slide, x, y, w, 0.4, [[R(label, GO, 11, accent, True, 120)]])
        hline(slide, x, y + 0.42, min(w, 2.2), prof["rule"], 1.0)
    else:  # plain
        text(slide, x, y, w, 0.4, [[R(label, GO, 11, accent, True, 100)]])


# ---- cover レイアウト微調整(bgsig ごとの左右グラフィックとの衝突回避) ----
_COVER_LAYOUT = {
    "rail":       dict(x=2.3, w=9.4, align=L, kx=2.3, kw=2.2),
    "diagonal":   dict(x=0.9, w=7.0, align=L, kx=0.9, kw=2.6),
    "lines":      dict(x=0.9, w=7.2, align=L, kx=0.9, kw=2.6),
    "circle":     dict(x=0.86, w=8.4, align=L, kx=0.9, kw=2.0),
    "colorblock": dict(x=0.88, w=7.4, align=L, kx=0.9, kw=2.4),
    "frame":      dict(x=1.5, w=10.3, align=C, kx=1.5, kw=10.3),
    "airyrule":   dict(x=4.95, w=7.6, align=L, kx=0.9, kw=3.4),
    "spine":      dict(x=4.95, w=7.5, align=L, kx=0.9, kw=3.4),
    "band":       dict(x=0.9, w=11.5, align=L, kx=0.9, kw=2.6),
    "panel":      dict(x=0.88, w=8.0, align=L, kx=0.9, kw=2.8),
    "glassdark":  dict(x=0.9, w=9.0, align=L, kx=0.9, kw=2.9),
    "horizon":    dict(x=0.88, w=9.0, align=L, kx=0.9, kw=2.9),
    "ghostnum":   dict(x=0.88, w=8.0, align=L, kx=0.9, kw=2.0),
    "wave":       dict(x=0.88, w=9.0, align=L, kx=0.9, kw=2.4),
    "bottomband": dict(x=0.88, w=8.0, align=L, kx=0.9, kw=2.3),
    "mono":       dict(x=0.9, w=11.6, align=L, kx=0.9, kw=2.9),
    "_default":   dict(x=0.9, w=11.5, align=L, kx=0.9, kw=2.6),
}


def _hero_stat(slide, prof, x, y, w, h, stat):
    """cover 用のヒーロー統計カード。meta["stat"] が指定された場合のみ描画する。"""
    card_plain(slide, x, y, w, h, prof["card"], shadow=prof.get("shadow"))
    ty = y + 0.4
    label = stat.get("label")
    if label:
        text(slide, x + 0.35, ty, w - 0.7, 0.4, [[R(str(label), prof["body_font"], 13, prof["muted"])]])
        ty += 0.5
    value_color = prof["ink"] if prof.get("dark") else prof["accent"]
    text(slide, x + 0.3, ty, w - 0.6, h * 0.4,
         [[R(str(stat.get("value", "")), prof["heading_font"], 48, value_color, True)]])
    note = stat.get("note")
    if note:
        pill(slide, x + 0.35, y + h - 0.65, min(w - 0.7, 2.2), 0.48,
             str(note), prof["accent"], prof.get("on_accent", "FFFFFF"), 11)


def render_cover(prs, prof, meta):
    """cover スライドを1枚描画して返す。

    meta: title(必須) / subtitle / date / author / kicker / audience / stat を利用する。
    """
    meta = meta or {}
    s = slide(prs, prof["bg"])
    _bg(s, prof, "cover")
    layout = _COVER_LAYOUT.get(prof.get("bgsig"), _COVER_LAYOUT["_default"])
    align = layout["align"]
    kicker_label = meta.get("kicker") or meta.get("audience")
    if kicker_label:
        _kicker(s, prof, layout["kx"], 1.3, layout["kw"], 0.48, str(kicker_label))
    title = meta.get("title", "")
    text(s, layout["x"], 2.55, layout["w"], 2.1,
         [[R(str(title), prof["heading_font"], 46, prof["ink"], True)]], align=align, ls=1.05)
    y = 4.85
    subtitle = meta.get("subtitle")
    if subtitle:
        text(s, layout["x"], y, layout["w"], 0.55,
             [[R(str(subtitle), prof["body_font"], 17, prof["muted"])]], align=align)
        y += 0.62
    bits = [str(b) for b in (meta.get("date"), meta.get("author")) if b]
    if bits:
        text(s, layout["x"], y, layout["w"], 0.4,
             [[R("　".join(bits), prof["body_font"], 13, prof["muted"])]], align=align)
    stat = meta.get("stat")
    if stat:
        _hero_stat(s, prof, 8.75, 2.15, 3.7, 3.15, stat)
    return s


def render_section(prs, prof, number, title):
    """section(中扉)スライドを1枚描画して返す。number は "03" のような文字列。"""
    s = slide(prs, prof["bg"])
    _bg(s, prof, "section")
    bgsig = prof.get("bgsig")
    ink = prof["ink"]
    accent = prof["accent"]
    on_accent = prof.get("on_accent") or "FFFFFF"
    label = f"SECTION {number}" if number else "SECTION"
    heading = prof["heading_font"]
    if bgsig in ("panel", "band", "circle"):
        _kicker(s, prof, 0.95, 2.35, 2.0, 0.5, label)
        text(s, 0.9, 2.9, 11.5, 1.8, [[R(str(title), heading, 48, on_accent, True)]], anchor=MID)
    elif bgsig in ("glassdark", "horizon"):
        _kicker(s, prof, 0.95, 2.35, 2.0, 0.5, label)
        text(s, 0.9, 3.0, 11.5, 1.6, [[R(str(title), heading, 48, ink, True)]], anchor=MID)
    elif bgsig == "rail":
        text(s, 0.6, 3.0, 4.6, 1.4,
             [[R(str(number), GO, 18, on_accent, True)], [R(str(title), heading, 30, on_accent, True)]],
             ls=1.3, anchor=MID)
    elif bgsig == "lines":
        text(s, 0.9, 3.1, 4.0, 1.4,
             [[R(str(number), GO, 20, on_accent, True)], [R(str(title), heading, 32, on_accent, True)]],
             ls=1.3, anchor=MID)
    elif bgsig == "diagonal":
        text(s, 0.9, 3.3, 3.0, 0.5, [[R(str(number), GO, 20, accent, True)]])
        text(s, 7.2, 3.0, 5.6, 1.4, [[R(str(title), heading, 34, on_accent, True)]], ls=1.1)
    elif bgsig == "colorblock":
        # 番号はタイトル同様アクセント全面ブロック(y<3.4)の中に収め、ブロック色と
        # 同化しないよう on_accent で描く(旧: accent 色でブロックに埋没していた)。
        text(s, 0.9, 1.1, 11.5, 1.4, [[R(str(title), heading, 40, on_accent, True)]])
        text(s, 0.9, 2.9, 2.0, 0.4, [[R(str(number), GO, 16, on_accent, True)]])
    elif bgsig == "spine":
        text(s, 0.9, 3.0, 5.4, 1.4,
             [[R(str(number), GO, 14, prof["muted"], False, 60)], [R(str(title), heading, 34, ink, True)]],
             ls=1.3, anchor=MID)
    else:
        _kicker(s, prof, 0.95, 2.35, 2.0, 0.5, label)
        text(s, 0.9, 3.0, 11.5, 1.4, [[R(str(title), heading, 44, ink, True)]], align=C, anchor=MID)
    return s


# ---- body(message/stats) レイアウト微調整(bgsig ごとの背景装飾との衝突回避) ----
# キー: x/w=本文の左端・幅、ky=キッカーの y、hy=見出しの y、cy=本文コンテンツの開始 y、
# cbottom=本文コンテンツが収まる下限 y。各値は _bg(variant="body") が描く装飾
# (帯・レール・波・フレーム・スパイン等)を避けるよう調整してある。
_BODY_LAYOUT = {
    "band":       dict(x=0.9, w=11.5, ky=1.05, hy=1.65, cy=2.6, cbottom=7.1),
    "rail":       dict(x=2.3, w=9.4,  ky=0.7,  hy=1.3,  cy=2.5, cbottom=7.1),
    "wave":       dict(x=0.9, w=11.5, ky=0.6,  hy=1.25, cy=2.5, cbottom=6.3),
    "bottomband": dict(x=0.9, w=11.5, ky=0.6,  hy=1.25, cy=2.5, cbottom=6.65),
    "frame":      dict(x=1.3, w=10.7, ky=0.85, hy=1.4,  cy=2.7, cbottom=6.7),
    "spine":      dict(x=1.2, w=11.2, ky=0.6,  hy=0.85, cy=2.2, cbottom=6.9),
    "lines":      dict(x=0.9, w=11.5, ky=0.6,  hy=1.2,  cy=2.6, cbottom=7.1),
    "_default":   dict(x=0.9, w=11.5, ky=0.6,  hy=1.25, cy=2.5, cbottom=7.1),
}


def _body_layout(prof):
    return _BODY_LAYOUT.get(prof.get("bgsig"), _BODY_LAYOUT["_default"])


# ---- カード描画(prof["card_style"] に応じたバリエーション) ----
def _card(s, prof, x, y, w, h):
    style = prof.get("card_style", "plain")
    fill = prof.get("card") or prof["bg"]
    shadow = prof.get("shadow")
    accent = prof["accent"]
    if style == "topstrip":
        card_topstrip(s, x, y, w, h, fill, accent, shadow=shadow)
    elif style == "leftbar":
        card_leftbar(s, x, y, w, h, fill, accent, shadow=shadow)
    elif style == "outline":
        card_outline(s, x, y, w, h, accent, lw=1.25, fill=fill)
    elif style == "rounded":
        card_plain(s, x, y, w, h, fill, shadow=shadow, radius=0.09)
    elif style == "glass":
        line = prof.get("glow") or accent
        card_plain(s, x, y, w, h, fill, line=line, lw=1.0, lalpha=40)
    else:  # plain
        card_plain(s, x, y, w, h, fill, shadow=shadow)


def _stat_block(s, prof, x, y, w, h, item, value_size, label_size, note_mode="pill"):
    """カード(またはコンテンツ枠)内に label/value/note を積み上げて描く。"""
    accent = prof["accent"]
    ty = y
    label = item.get("label")
    if label:
        text(s, x, ty, w, 0.34, [[R(str(label), prof["body_font"], label_size, prof["muted"])]])
        ty += 0.38
    vh = min(value_size / 72.0 * 1.35, max(h - (ty - y) - 0.4, 0.5))
    text(s, x, ty, w, vh, [[R(str(item.get("value", "")), prof["heading_font"], value_size, accent, True)]])
    ty += vh
    note = item.get("note")
    if note:
        if note_mode == "pill" and w >= 1.8 and (y + h - ty) >= 0.4:
            ny = min(ty + 0.05, y + h - 0.44)
            pill(s, x, ny, min(w, 2.2), 0.4, str(note), accent, prof.get("on_accent", "FFFFFF"), 10)
        else:
            text(s, x, ty + 0.03, w, 0.32, [[R(str(note), prof["body_font"], max(label_size - 1, 9), accent, True)]])


# ---- stats の4アーキタイプ ----
def _stats_bento(s, prof, area, items):
    """藍/藍鉄/墨/石板/鉄紺 — 先頭 item を大型ヒーローカード、残りを右側の小カード群に。"""
    x, y, w, bottom = area["x"], area["y"], area["w"], area["bottom"]
    h = max(bottom - y, 1.0)
    hero, rest = items[0], items[1:]
    hero_w = w * 0.52
    _card(s, prof, x, y, hero_w, h)
    _stat_block(s, prof, x + 0.4, y + 0.4, hero_w - 0.8, h - 0.8, hero, value_size=64, label_size=14)
    if not rest:
        return
    rx = x + hero_w + 0.25
    rw = w - hero_w - 0.25
    n = len(rest)
    gap = 0.22
    rh = max((h - gap * (n - 1)) / n, 0.7)
    vsize = 34 if n <= 2 else (24 if n <= 4 else 17)
    ry = y
    for it in rest:
        _card(s, prof, rx, ry, rw, rh)
        _stat_block(s, prof, rx + 0.32, ry + 0.28, rw - 0.64, rh - 0.56, it, value_size=vsize, label_size=12)
        ry += rh + gap


def _stats_cards(s, prof, area, items):
    """常磐/青碧/山吹/彩層/桜/亜麻 — 均等カード。4枚を超えたら2行に折り返す。"""
    x, y, w, bottom = area["x"], area["y"], area["w"], area["bottom"]
    n = len(items)
    cols = min(n, 4)
    rows = math.ceil(n / cols)
    gap = 0.22
    total_h = max(bottom - y, 1.0)
    row_h = max((total_h - gap * (rows - 1)) / rows, 1.2)
    col_w = (w - gap * (cols - 1)) / cols
    vsize = 40 if n <= 3 else (28 if n <= 6 else 20)
    for i, it in enumerate(items):
        r, c = divmod(i, cols)
        cx = x + c * (col_w + gap)
        cy = y + r * (row_h + gap)
        _card(s, prof, cx, cy, col_w, row_h)
        _stat_block(s, prof, cx + 0.3, cy + 0.3, col_w - 0.6, row_h - 0.6, it, value_size=vsize, label_size=13)


def _stats_poster(s, prof, area, items):
    """白磁/朱 — 単一 item は超特大 value、複数は大きな数字を縦に並べる。"""
    x, y, w, bottom = area["x"], area["y"], area["w"], area["bottom"]
    if len(items) == 1:
        it = items[0]
        text(s, x, y, w, 2.4, [[R(str(it.get("value", "")), prof["heading_font"], 140, prof["accent"], True)]])
        ty = y + 2.5
        label = it.get("label")
        if label:
            text(s, x, ty, w, 0.5, [[R(str(label), prof["body_font"], 20, prof["ink"], True)]])
            ty += 0.55
        note = it.get("note")
        if note:
            text(s, x, ty, w, 0.5, [[R(str(note), prof["body_font"], 15, prof["muted"])]])
        return
    n = len(items)
    row_h = max((bottom - y) / n, 0.7)
    vsize = 44 if n <= 3 else (30 if n <= 6 else 22)
    for i, it in enumerate(items):
        ry = y + i * row_h
        text(s, x, ry, 2.4, row_h,
             [[R(str(it.get("value", "")), prof["heading_font"], vsize, prof["accent"], True)]], anchor=MID)
        label_note = str(it.get("label", ""))
        note = it.get("note")
        if note:
            label_note += f"　{note}"
        text(s, x + 2.6, ry, w - 2.6, row_h, [[R(label_note, prof["body_font"], 16, prof["ink"])]], anchor=MID)


def _stats_list(s, prof, area, items):
    """霞/藤/明朝 — value を大きく、label/note を行として並べる。"""
    x, y, w, bottom = area["x"], area["y"], area["w"], area["bottom"]
    n = len(items)
    row_h = max((bottom - y) / n, 0.65)
    vsize = 32 if n <= 4 else (22 if n <= 6 else 16)
    for i, it in enumerate(items):
        ry = y + i * row_h
        hline(s, x, ry, w, prof["rule"], 1.0)
        text(s, x, ry + 0.1, 2.6, row_h - 0.2,
             [[R(str(it.get("value", "")), prof["heading_font"], vsize, prof["accent"], True)]], anchor=MID)
        lines = [[R(str(it.get("label", "")), prof["body_font"], 15, prof["ink"], True)]]
        note = it.get("note")
        if note:
            lines.append([R(str(note), prof["body_font"], 12, prof["muted"])])
        text(s, x + 2.8, ry + 0.1, w - 2.8, row_h - 0.2, lines, ls=1.2, anchor=MID)


_STATS_DISPATCH = {
    "bento": _stats_bento, "cards": _stats_cards, "poster": _stats_poster, "list": _stats_list,
}


# ---- message の箇条書きマーカー(prof["bullet"]) ----
def _bullet_marker(s, prof, x, y, w, row_h, txt, font_size, idx):
    bullet = prof.get("bullet", "square")
    accent = prof["accent"]
    ink = prof["ink"]
    bf = prof["body_font"]
    mw = 0.55
    my = y + row_h * 0.5
    if bullet == "square":
        rect(s, x, my - 0.06, 0.13, 0.13, fill=accent)
    elif bullet == "dash":
        text(s, x, y, mw, row_h, [[R("—", prof["heading_font"], font_size, accent, True)]], anchor=MID)
    elif bullet == "chevron":
        poly(s, [(x, my - 0.09), (x + 0.16, my), (x, my + 0.09)], fill=accent)
    elif bullet == "ring":
        oval(s, x, my - 0.09, 0.18, 0.18, line=accent, lw=1.5)
    elif bullet == "number":
        text(s, x, y, mw + 0.15, row_h, [[R(str(idx), prof["heading_font"], font_size, accent, True)]], anchor=MID)
    else:  # tick
        rect(s, x, my - 0.012, 0.22, 0.03, fill=prof.get("rule") or accent)
    text(s, x + mw + 0.15, y, w - mw - 0.15, row_h, [[R(txt, bf, font_size, ink)]], anchor=MID, ls=1.15)


def render_message(prs, prof, headline, body):
    """message スライドを1枚描画して返す。body は文字列のリスト(空/None安全)。"""
    s = slide(prs, prof["bg"])
    _bg(s, prof, "body")
    lay = _body_layout(prof)
    _kicker(s, prof, lay["x"], lay["ky"], 2.2, 0.42, "POINT")
    text(s, lay["x"], lay["hy"], lay["w"], 1.0,
         [[R(str(headline), prof["heading_font"], 30, prof["ink"], True)]], ls=1.05)
    items = [str(b) for b in (body or []) if b is not None and str(b) != ""]
    if not items:
        return s
    n = len(items)
    top, bottom = lay["cy"], lay["cbottom"]
    row_h = max((bottom - top) / n, 0.45)
    font_size = 18 if n <= 3 else (16 if n <= 5 else 13)
    y = top
    for i, it in enumerate(items, start=1):
        _bullet_marker(s, prof, lay["x"], y, lay["w"], row_h, it, font_size, i)
        y += row_h
    return s


def render_stats(prs, prof, headline, items):
    """stats スライドを1枚描画して返す。items は {value,label,note?} のリスト(空安全)。"""
    s = slide(prs, prof["bg"])
    _bg(s, prof, "body")
    lay = _body_layout(prof)
    _kicker(s, prof, lay["x"], lay["ky"], 2.2, 0.42, "DATA")
    text(s, lay["x"], lay["hy"], lay["w"], 1.0,
         [[R(str(headline), prof["heading_font"], 30, prof["ink"], True)]], ls=1.05)
    clean = [it for it in (items or []) if isinstance(it, dict)]
    if not clean:
        return s
    area = dict(x=lay["x"], y=lay["cy"], w=lay["w"], bottom=lay["cbottom"])
    fn = _STATS_DISPATCH.get(prof.get("layout"), _stats_cards)
    fn(s, prof, area, clean)
    return s
