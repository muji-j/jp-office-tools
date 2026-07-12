from pathlib import Path
import pytest
import jp_slides
pptx = pytest.importorskip("pptx")
from pptx import Presentation


def _render(tmp_path, slide):
    c = jp_slides.parse_content({"meta": {"title": "t"}, "theme": "藍",
                                 "slides": [{"type": "cover", "title": "t"}, slide]})
    return Presentation(jp_slides.render_deck(c, out=str(tmp_path / "o.pptx")))


def test_table_slide(tmp_path):
    prs = _render(tmp_path, {"type": "table", "headline": "比較表",
                             "columns": ["評価軸", "A案", "B案"],
                             "rows": [["コスト", "低", "高"], ["納期", "3週", "5週"]]})
    tbls = [sh for sh in prs.slides[1].shapes if sh.has_table]
    assert tbls, "表が無い"
    t = tbls[0].table
    assert len(t.rows) == 3 and len(t.columns) == 3
    assert t.cell(0, 0).text == "評価軸"
    assert t.cell(2, 1).text == "3週"


def test_image_slide(tmp_path):
    # 1x1 PNG フィクスチャ
    png = tmp_path / "x.png"
    png.write_bytes(bytes.fromhex(
        "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c4"
        "890000000d4944415478da6360000002000154a24f9f0000000049454e44ae426082"))
    prs = _render(tmp_path, {"type": "image", "headline": "推移", "image": str(png), "caption": "図1"})
    pics = [sh for sh in prs.slides[1].shapes if sh.shape_type == 13]  # PICTURE
    assert pics, "画像が挿入されていない"


def test_image_missing_path_warns_but_builds(tmp_path, capsys):
    prs = _render(tmp_path, {"type": "image", "headline": "無い画像",
                             "image": str(tmp_path / "nope.png"), "caption": "図2"})
    err = capsys.readouterr().err
    assert "画像" in err  # 警告が出る
    assert len(prs.slides) == 2  # デッキは生成される


def test_section_slide(tmp_path):
    prs = _render(tmp_path, {"type": "section", "title": "第2部 詳細データ"})
    txt = "\n".join(sh.text_frame.text for sh in prs.slides[1].shapes if sh.has_text_frame)
    assert "第2部 詳細データ" in txt


def test_table_slide_empty_cell_no_crash(tmp_path):
    prs = _render(tmp_path, {"type": "table", "headline": "備考表",
                             "columns": ["項目", "値", "備考"],
                             "rows": [["備考", "", "特記なし"]]})
    tbls = [sh for sh in prs.slides[1].shapes if sh.has_table]
    assert tbls, "表が無い"
    t = tbls[0].table
    assert t.cell(1, 1).text == ""


def test_table_slide_short_row_pads_empty(tmp_path):
    prs = _render(tmp_path, {"type": "table", "headline": "比較表",
                             "columns": ["軸", "A案", "B案"],
                             "rows": [["コスト", "低"]]})
    tbls = [sh for sh in prs.slides[1].shapes if sh.has_table]
    assert tbls, "表が無い"
    t = tbls[0].table
    assert t.cell(1, 0).text == "コスト"
    assert t.cell(1, 1).text == "低"
    assert t.cell(1, 2).text == ""


def test_table_slide_empty_header_no_crash(tmp_path):
    prs = _render(tmp_path, {"type": "table", "headline": "空ヘッダー",
                             "columns": ["項目", ""],
                             "rows": [["A", "B"]]})
    tbls = [sh for sh in prs.slides[1].shapes if sh.has_table]
    assert tbls, "表が無い"
    t = tbls[0].table
    assert t.cell(0, 1).text == ""


def test_table_slide_many_rows_height_capped(tmp_path):
    rows = [[f"項目{i}", f"値{i}"] for i in range(12)]
    prs = _render(tmp_path, {"type": "table", "headline": "長い表",
                             "columns": ["項目", "値"], "rows": rows})
    tbls = [sh for sh in prs.slides[1].shapes if sh.has_table]
    assert tbls, "表が無い"
    shp = tbls[0]
    # スライド高(7.5インチ = 6858000 EMU)を超えない
    assert shp.top + shp.height <= 6858000


def test_image_slide_tall_image_fits_and_caption_below(tmp_path):
    from PIL import Image
    png = tmp_path / "tall.png"
    Image.new("RGB", (100, 300)).save(png)
    prs = _render(tmp_path, {"type": "image", "headline": "縦長画像",
                             "image": str(png), "caption": "図3"})
    slide = prs.slides[1]
    pics = [sh for sh in slide.shapes if sh.shape_type == 13]  # PICTURE
    assert pics, "画像が挿入されていない"
    pic = pics[0]
    slide_height = 6858000  # 7.5in in EMU
    assert pic.top + pic.height <= slide_height
    # キャプションのテキストボックスを見つける(画像より下にあり、スライド内に収まる)
    textboxes = [sh for sh in slide.shapes
                 if sh.has_text_frame and "図3" in sh.text_frame.text]
    assert textboxes, "キャプションが見つからない"
    cap = textboxes[0]
    assert cap.top < slide_height
    assert cap.top >= pic.top + pic.height
