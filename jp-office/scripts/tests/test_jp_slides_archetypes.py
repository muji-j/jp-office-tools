from pathlib import Path
import pytest
import jp_slides
pptx = pytest.importorskip("pptx")
from pptx import Presentation
from pptx.dml.color import RGBColor


def _deck(theme):
    return jp_slides.parse_content({
        "meta": {"title": "検証"},
        "theme": theme,
        "slides": [
            {"type": "cover", "title": "検証資料"},
            {"type": "message", "headline": "結論を一文で言い切る", "body": ["根拠A", "根拠B"]},
            {"type": "message", "headline": "二枚目の主張"},
        ],
    })


@pytest.mark.parametrize("theme", list(jp_slides.THEMES))
def test_every_theme_builds(tmp_path, theme):
    out = jp_slides.render_deck(_deck(theme), out=str(tmp_path / f"{theme}.pptx"))
    prs = Presentation(out)
    assert len(prs.slides) == 3


@pytest.mark.parametrize("theme", ["墨", "藍鉄"])
def test_dark_theme_background(tmp_path, theme):
    out = jp_slides.render_deck(_deck(theme), out=str(tmp_path / "d.pptx"))
    prs = Presentation(out)
    bg = prs.slides[0].background.fill.fore_color.rgb
    assert bg == RGBColor.from_string(jp_slides.THEMES[theme].bg)


def test_header_band_present(tmp_path):
    # 鉄紺(header-band): accent #1B2A4A 帯シェイプが存在する
    out = jp_slides.render_deck(_deck("鉄紺"), out=str(tmp_path / "h.pptx"))
    prs = Presentation(out)
    band = False
    for sh in prs.slides[1].shapes:
        try:
            if sh.fill.fore_color.rgb == RGBColor.from_string("1B2A4A"):
                band = True
        except Exception:
            pass
    assert band


def test_color_block_two_colors(tmp_path):
    # 彩層(color-block): accent #2D3E78 + accent2 #F0A500 の両方が存在する
    out = jp_slides.render_deck(_deck("彩層"), out=str(tmp_path / "c.pptx"))
    prs = Presentation(out)
    seen = set()
    for sh in prs.slides[0].shapes:
        try:
            seen.add(str(sh.fill.fore_color.rgb))
        except Exception:
            pass
    assert "2D3E78" in seen and "F0A500" in seen
