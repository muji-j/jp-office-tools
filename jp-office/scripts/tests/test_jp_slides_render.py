import json
from pathlib import Path
import pytest
import jp_slides

pptx = pytest.importorskip("pptx")  # python-pptx 未導入なら skip
from pptx import Presentation
from pptx.util import Emu


def _content(theme="藍"):
    return {
        "meta": {"title": "月次営業報告", "date": "2026-07-13", "author": "山田太郎"},
        "theme": theme,
        "pattern": "conclusion",
        "slides": [
            {"type": "cover", "title": "月次営業報告", "subtitle": "2026年6月度",
             "date": "2026-07-13", "author": "山田太郎"},
            {"type": "message", "headline": "6月の売上は前月比12%増加した",
             "body": ["新規顧客が8件増加", "既存顧客の解約はゼロ"]},
        ],
    }


def _all_text(prs):
    out = []
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame:
                out.append(sh.text_frame.text)
    return "\n".join(out)


def test_render_creates_valid_pptx(tmp_path):
    c = jp_slides.parse_content(_content("藍"))
    out = jp_slides.render_deck(c, out=str(tmp_path / "deck.pptx"))
    p = Path(out)
    assert p.exists() and p.suffix == ".pptx" and p.stat().st_size > 0
    prs = Presentation(out)
    assert len(prs.slides) == 2


def test_render_contains_title_and_headline(tmp_path):
    c = jp_slides.parse_content(_content("藍"))
    out = jp_slides.render_deck(c, out=str(tmp_path / "deck.pptx"))
    text = _all_text(Presentation(out))
    assert "月次営業報告" in text
    assert "6月の売上は前月比12%増加した" in text


def test_render_slide_size_16by9(tmp_path):
    c = jp_slides.parse_content(_content("藍"))
    out = jp_slides.render_deck(c, out=str(tmp_path / "deck.pptx"))
    prs = Presentation(out)
    # 13.333in x 7.5in ≈ 12192000 x 6858000 EMU
    assert abs(prs.slide_width - Emu(12192000)) < Emu(20000)
    assert abs(prs.slide_height - Emu(6858000)) < Emu(20000)


def test_render_accent_color_present(tmp_path):
    # accent-bar テーマの藍: accent #2A4B8D の塗り図形が存在
    c = jp_slides.parse_content(_content("藍"))
    out = jp_slides.render_deck(c, out=str(tmp_path / "deck.pptx"))
    prs = Presentation(out)
    found = False
    for s in prs.slides:
        for sh in s.shapes:
            try:
                if sh.fill.type is not None and sh.fill.fore_color.rgb == pptx.dml.color.RGBColor.from_string("2A4B8D"):
                    found = True
            except Exception:
                pass
    assert found, "accent色の図形が見つからない"


def test_build_cli(tmp_path, capsys):
    cpath = tmp_path / "c.json"
    cpath.write_text(json.dumps(_content("霞"), ensure_ascii=False), encoding="utf-8")
    outp = tmp_path / "o.pptx"
    rc = jp_slides.main(["jp_slides.py", "build", str(cpath), "--out", str(outp)])
    assert rc == 0
    assert outp.exists()
    printed = capsys.readouterr().out.strip()
    assert str(outp) in printed


def test_build_cli_theme_override(tmp_path):
    cpath = tmp_path / "c.json"
    body = _content("藍")
    cpath.write_text(json.dumps(body, ensure_ascii=False), encoding="utf-8")
    outp = tmp_path / "o.pptx"
    jp_slides.main(["jp_slides.py", "build", str(cpath), "--theme", "霞", "--out", str(outp)])
    # 霞 bg #FAFAFA が背景に反映(藍のFFFFFFではない)
    prs = Presentation(str(outp))
    bg = prs.slides[0].background.fill.fore_color.rgb
    assert bg == pptx.dml.color.RGBColor.from_string("FAFAFA")
