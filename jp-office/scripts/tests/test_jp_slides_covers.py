import pytest

pytest.importorskip("pptx")

import jp_slides
import jp_slides_design as D
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

THEME_KEYS = list(jp_slides.THEME_PROFILES)


def _all_text(slide):
    chunks = []
    for shp in slide.shapes:
        if shp.has_text_frame:
            chunks.append(shp.text_frame.text)
    return "\n".join(chunks)


def _shape_count(slide):
    return len(slide.shapes)


@pytest.mark.parametrize("key", THEME_KEYS)
def test_cover_and_section_build_valid_pptx(tmp_path, key):
    prof = jp_slides.THEME_PROFILES[key]
    prs = D.new_prs()
    D.render_cover(prs, prof, {"title": "事業戦略アップデート", "date": "2026年7月", "author": "山田太郎"})
    D.render_section(prs, prof, "03", "下期の重点施策")
    out = tmp_path / f"{key}.pptx"
    prs.save(str(out))

    reopened = Presentation(str(out))
    assert len(reopened.slides) == 2


@pytest.mark.parametrize("key", THEME_KEYS)
def test_cover_title_text_present(tmp_path, key):
    prof = jp_slides.THEME_PROFILES[key]
    prs = D.new_prs()
    cover = D.render_cover(prs, prof, {"title": "事業戦略アップデート", "date": "2026年7月", "author": "山田太郎"})
    assert "事業戦略アップデート" in _all_text(cover)


@pytest.mark.parametrize("key", THEME_KEYS)
def test_section_title_and_number_present(tmp_path, key):
    prof = jp_slides.THEME_PROFILES[key]
    prs = D.new_prs()
    section = D.render_section(prs, prof, "03", "下期の重点施策")
    text = _all_text(section)
    assert "下期の重点施策" in text
    assert "03" in text


@pytest.mark.parametrize("key", ["墨", "藍鉄"])
def test_dark_theme_background_is_dark_bg(key):
    prof = jp_slides.THEME_PROFILES[key]
    prs = D.new_prs()
    cover = D.render_cover(prs, prof, {"title": "タイトル"})
    fill_rgb = cover.background.fill.fore_color.rgb
    assert fill_rgb == RGBColor.from_string(prof["bg"])
    section = D.render_section(prs, prof, "01", "セクション")
    assert section.background.fill.fore_color.rgb == RGBColor.from_string(prof["bg"])


def test_diagonal_theme_has_freeform_wedge():
    prof = jp_slides.THEME_PROFILES["常磐"]
    prs = D.new_prs()
    cover = D.render_cover(prs, prof, {"title": "タイトル"})
    shape_types = [shp.shape_type for shp in cover.shapes]
    assert MSO_SHAPE_TYPE.FREEFORM in shape_types


def test_sakura_theme_has_wave_freeform():
    prof = jp_slides.THEME_PROFILES["桜"]
    prs = D.new_prs()
    cover = D.render_cover(prs, prof, {"title": "タイトル"})
    shape_types = [shp.shape_type for shp in cover.shapes]
    assert shape_types.count(MSO_SHAPE_TYPE.FREEFORM) >= 2


def test_seiheki_theme_has_split_line_field():
    prof = jp_slides.THEME_PROFILES["青碧"]
    prs = D.new_prs()
    baseline_prs = D.new_prs()
    baseline_prof = jp_slides.THEME_PROFILES["白磁"]  # mono: 装飾図形なし
    baseline_cover = D.render_cover(baseline_prs, baseline_prof, {"title": "タイトル"})
    cover = D.render_cover(prs, prof, {"title": "タイトル"})
    # seiheki(青碧) は右側フィールド+境界線+分割線で、装飾なしテーマより図形数が多い
    assert _shape_count(cover) > _shape_count(baseline_cover)


def test_mono_theme_has_no_extra_background_shapes():
    prof = jp_slides.THEME_PROFILES["白磁"]
    prs = D.new_prs()
    cover = D.render_cover(prs, prof, {"title": "タイトル"})
    # mono は kicker/title のテキストボックスのみ(kicker が無ければタイトルのみ)
    for shp in cover.shapes:
        assert shp.shape_type != MSO_SHAPE_TYPE.FREEFORM


def test_unknown_bgsig_raises():
    with pytest.raises(ValueError):
        D._bg(None, {"bgsig": "no-such-signature"}, "cover")


def test_render_cover_handles_missing_optional_meta():
    prof = jp_slides.THEME_PROFILES["藍"]
    prs = D.new_prs()
    cover = D.render_cover(prs, prof, {"title": "最小構成"})
    assert "最小構成" in _all_text(cover)


def _run_colors(slide):
    """スライド上の全テキスト run について (文字列, RGBColor) のリストを返す。"""
    out = []
    for shp in slide.shapes:
        if not shp.has_text_frame:
            continue
        for p in shp.text_frame.paragraphs:
            for r in p.runs:
                out.append((r.text, r.font.color.rgb))
    return out


def test_saisou_section_number_color_differs_from_block_fill():
    # 彩層(colorblock) の section は上部にアクセント全面ブロックを敷く。
    # 番号テキストがそのブロック色と同色だと埋没するため、異なる色であることを検証する。
    prof = jp_slides.THEME_PROFILES["彩層"]
    prs = D.new_prs()
    section = D.render_section(prs, prof, "03", "下期の重点施策")
    block_color = RGBColor.from_string(prof["accent"])
    number_colors = [c for (t, c) in _run_colors(section) if t == "03"]
    assert number_colors, "セクション番号のテキストランが見つかりません。"
    for c in number_colors:
        assert c != block_color


@pytest.mark.parametrize("variant", ["cover", "body"])
def test_ai_panel_has_no_tint_side_panel(variant):
    # 藍(ai) の原典は cover/body に背景装飾を持たない。派生の淡色ティントパネル
    # (右側の _tint(accent) 矩形)が復活していないことを回帰的に確認する。
    prof = jp_slides.THEME_PROFILES["藍"]
    prs = D.new_prs()
    s = D.slide(prs, prof["bg"])
    D._bg(s, prof, variant)
    assert len(s.shapes) == 0


@pytest.mark.parametrize("variant", ["cover", "body"])
def test_sumi_glassdark_has_no_glow_oval(variant):
    # 墨(sumi) の原典は cover/body にヘアライン・グロー円を持たない。
    prof = jp_slides.THEME_PROFILES["墨"]
    prs = D.new_prs()
    s = D.slide(prs, prof["bg"])
    D._bg(s, prof, variant)
    assert len(s.shapes) == 0


def test_render_cover_renders_hero_stat_when_present():
    prof = jp_slides.THEME_PROFILES["藍"]
    prs = D.new_prs()
    cover = D.render_cover(prs, prof, {
        "title": "事業戦略アップデート",
        "stat": {"value": "303億", "label": "上期 売上高", "note": "前年比 +14%"},
    })
    text = _all_text(cover)
    assert "303億" in text
    assert "上期 売上高" in text
    assert "前年比 +14%" in text


# ---- hero-stat ガード回帰(T2 由来の _HERO_STAT_BLOCKED_BGSIGS) ----
# 右側領域(8.75〜12.45in)がシグネチャ背景の右側装飾と衝突する6テーマでは、
# stat が meta に指定されてもヒーローカードを描画してはならない。
_HERO_STAT_BLOCKED_THEMES = ["常磐", "青碧", "朱", "彩層", "鉄紺", "石板"]


@pytest.mark.parametrize("key", _HERO_STAT_BLOCKED_THEMES)
def test_hero_stat_blocked_for_conflicting_bgsig_themes(key):
    prof = jp_slides.THEME_PROFILES[key]
    assert prof["bgsig"] in D._HERO_STAT_BLOCKED_BGSIGS, (
        f"{key} の bgsig {prof['bgsig']!r} が _HERO_STAT_BLOCKED_BGSIGS から外れている"
    )
    prs = D.new_prs()
    cover = D.render_cover(prs, prof, {
        "title": "事業戦略アップデート",
        "stat": {"value": "303億", "label": "上期売上"},
    })
    text = _all_text(cover)
    assert "303億" not in text
    assert "上期売上" not in text


def test_hero_stat_rendered_for_non_blocked_bgsig_theme():
    # 藍(panel) は _HERO_STAT_BLOCKED_BGSIGS に含まれないため、対照としてヒーローカードが描画される。
    prof = jp_slides.THEME_PROFILES["藍"]
    assert prof["bgsig"] not in D._HERO_STAT_BLOCKED_BGSIGS
    prs = D.new_prs()
    cover = D.render_cover(prs, prof, {
        "title": "事業戦略アップデート",
        "stat": {"value": "303億", "label": "上期売上"},
    })
    text = _all_text(cover)
    assert "303億" in text
    assert "上期売上" in text
