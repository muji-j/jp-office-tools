from pathlib import Path
import pytest
import jp_slides
pptx = pytest.importorskip("pptx")
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor


def _content(**over):
    base = {"meta": {"title": "t"}, "theme": "藍",
            "slides": [{"type": "cover", "title": "表題"},
                       {"type": "message", "headline": "主張", "body": ["x"]}]}
    base.update(over)
    return base


def test_accent_override(tmp_path):
    c = jp_slides.parse_content(_content(accent="#E24A2B"))
    out = jp_slides.render_deck(c, out=str(tmp_path / "o.pptx"))
    prs = Presentation(out)
    seen = set()
    for s in prs.slides:
        for sh in s.shapes:
            try:
                seen.add(str(sh.fill.fore_color.rgb))
            except Exception:
                pass
    assert "E24A2B" in seen


def test_font_override(tmp_path):
    c = jp_slides.parse_content(_content(font="meiryo"))
    assert c["font"] == "meiryo"
    out = jp_slides.render_deck(c, out=str(tmp_path / "o.pptx"))
    # 開ければ可(フォント適用は _apply_font 経由)。回帰: 生成が壊れない
    assert len(Presentation(out).slides) == 2


def test_variant_dark_override(tmp_path):
    # ライトテーマ藍 + variant=dark → 背景が実際に暗色へリマップされる(旗だけでなく色も)
    c = jp_slides.parse_content(_content(theme="藍", variant="dark"))
    out = jp_slides.render_deck(c, out=str(tmp_path / "o.pptx"))
    prs = Presentation(out)
    bg = prs.slides[0].background.fill.fore_color.rgb
    assert bg == RGBColor.from_string("1C1C1E")
    assert len(prs.slides) == 2


def test_variant_light_on_dark_theme(tmp_path):
    # ダークテーマ墨 + variant=light → 背景が白へリマップ
    c = jp_slides.parse_content(_content(theme="墨", variant="light"))
    out = jp_slides.render_deck(c, out=str(tmp_path / "o.pptx"))
    bg = Presentation(str(tmp_path / "o.pptx")).slides[0].background.fill.fore_color.rgb
    assert bg == RGBColor.from_string("FFFFFF")


def test_template_mode_inherits_master(tmp_path):
    # フィクスチャ: 独自スライドサイズのテンプレを作り、そのサイズが継承されるか
    tpl = Presentation()
    tpl.slide_width = Inches(10)   # 4:3
    tpl.slide_height = Inches(7.5)
    tpath = tmp_path / "brand.pptx"
    tpl.save(str(tpath))
    c = jp_slides.parse_content(_content(template=str(tpath)))
    out = jp_slides.render_deck(c, out=str(tmp_path / "o.pptx"))
    prs = Presentation(out)
    assert abs(prs.slide_width - Inches(10)) < Inches(0.02)  # テンプレ寸法を継承
    assert len(prs.slides) >= 2
