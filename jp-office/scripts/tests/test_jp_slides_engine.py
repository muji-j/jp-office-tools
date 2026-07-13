import json

import pytest

import jp_slides

pptx = pytest.importorskip("pptx")
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches

THEME_KEYS = list(jp_slides.THEME_PROFILES)

_ONE_PX_PNG_HEX = (
    "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c4"
    "890000000d4944415478da6360000002000154a24f9f0000000049454e44ae426082"
)


def _all_text(prs):
    out = []
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame:
                out.append(sh.text_frame.text)
    return "\n".join(out)


def _full_content(theme="藍", **over):
    base = {
        "meta": {"title": "月次営業報告", "date": "2026年7月13日", "author": "山田太郎"},
        "theme": theme,
        "slides": [
            {"type": "cover", "title": "月次営業報告", "subtitle": "株式会社サンプル 2026年6月度",
             "date": "2026年7月13日", "author": "企画部 山田太郎"},
            {"type": "message", "headline": "6月の売上は前月比12%増加した",
             "body": ["新規契約が8件成立", "解約はゼロ"]},
            {"type": "stats", "headline": "主要指標の伸び", "items": [
                {"value": "182億", "label": "エネルギー事業", "note": "+18%"},
                {"value": "24億", "label": "SaaS事業"},
            ]},
            {"type": "table", "headline": "案の比較", "columns": ["評価軸", "A案", "B案"],
             "rows": [["コスト", "低", "高"], ["納期", "3週間", "5週間"]]},
            {"type": "image", "headline": "推移", "image": "", "caption": "図1"},
            {"type": "section", "title": "次期の重点施策", "number": "02"},
        ],
    }
    base.update(over)
    return base


# ---- build CLI: 全タイプ含むデッキ生成 ----

def test_build_cli_all_slide_types(tmp_path, capsys):
    content = _full_content()
    cpath = tmp_path / "c.json"
    cpath.write_text(json.dumps(content, ensure_ascii=False), encoding="utf-8")
    outp = tmp_path / "o.pptx"
    rc = jp_slides.main(["jp_slides.py", "build", str(cpath), "--out", str(outp)])
    assert rc == 0
    assert outp.exists()
    prs = Presentation(str(outp))
    assert len(prs.slides) == 6
    text = _all_text(prs)
    assert "月次営業報告" in text
    assert "6月の売上は前月比12%増加した" in text
    assert "182億" in text
    assert "次期の重点施策" in text
    tbls = [sh for sh in prs.slides[3].shapes if sh.has_table]
    assert tbls, "表が無い"
    assert tbls[0].table.cell(1, 0).text == "コスト"


# ---- 16テーマ スモーク(全タイプ) ----

@pytest.mark.parametrize("key", THEME_KEYS)
def test_render_deck_all_types_all_themes_smoke(tmp_path, key):
    content = jp_slides.parse_content(_full_content(theme=key))
    out = jp_slides.render_deck(content, out=str(tmp_path / f"{key}.pptx"))
    prs = Presentation(out)
    assert len(prs.slides) == 6


# ---- オーバーライド(accent/font/variant) ----

def test_accent_override_applied(tmp_path):
    content = jp_slides.parse_content(_full_content(theme="藍", accent="#E24A2B"))
    out = jp_slides.render_deck(content, out=str(tmp_path / "o.pptx"))
    prs = Presentation(out)
    seen = set()
    for s in prs.slides:
        for sh in s.shapes:
            try:
                seen.add(str(sh.fill.fore_color.rgb))
            except Exception:
                pass
    assert "E24A2B" in seen


def test_font_override_applied(tmp_path):
    content = jp_slides.parse_content(_full_content(theme="霞", font="meiryo"))
    out = jp_slides.render_deck(content, out=str(tmp_path / "o.pptx"))
    prs = Presentation(out)
    cover = prs.slides[0]
    names = [r.font.name for sh in cover.shapes if sh.has_text_frame
             for p in sh.text_frame.paragraphs for r in p.runs]
    assert "メイリオ" in names


def test_variant_dark_remaps_light_theme_background(tmp_path):
    content = jp_slides.parse_content(_full_content(theme="霞", variant="dark"))
    out = jp_slides.render_deck(content, out=str(tmp_path / "o.pptx"))
    prs = Presentation(out)
    bg = prs.slides[0].background.fill.fore_color.rgb
    assert bg == RGBColor.from_string("1C1C1E")


def test_variant_light_remaps_dark_theme_background(tmp_path):
    content = jp_slides.parse_content(_full_content(theme="墨", variant="light"))
    out = jp_slides.render_deck(content, out=str(tmp_path / "o.pptx"))
    prs = Presentation(out)
    bg = prs.slides[0].background.fill.fore_color.rgb
    assert bg == RGBColor.from_string("FFFFFF")


# ---- template モード ----

def test_template_mode_inherits_dimensions(tmp_path):
    tpl = Presentation()
    tpl.slide_width = Inches(10)  # 4:3
    tpl.slide_height = Inches(7.5)
    tpath = tmp_path / "brand.pptx"
    tpl.save(str(tpath))
    content = jp_slides.parse_content(_full_content(theme="藍", template=str(tpath)))
    out = jp_slides.render_deck(content, out=str(tmp_path / "o.pptx"))
    prs = Presentation(out)
    assert abs(prs.slide_width - Inches(10)) < Inches(0.02)
    assert len(prs.slides) == 6


# ---- 表: 空セル・短い行・行高上限(回帰) ----

def test_table_empty_cell_no_crash(tmp_path):
    content = jp_slides.parse_content({
        "theme": "藍",
        "slides": [{"type": "table", "headline": "備考表",
                    "columns": ["項目", "値", "備考"],
                    "rows": [["備考", "", "特記なし"]]}],
    })
    out = jp_slides.render_deck(content, out=str(tmp_path / "o.pptx"))
    prs = Presentation(out)
    tbls = [sh for sh in prs.slides[0].shapes if sh.has_table]
    assert tbls, "表が無い"
    t = tbls[0].table
    assert t.cell(1, 1).text == ""
    assert t.cell(1, 2).text == "特記なし"


def test_table_short_row_pads_empty(tmp_path):
    content = jp_slides.parse_content({
        "theme": "藍",
        "slides": [{"type": "table", "headline": "比較表",
                    "columns": ["軸", "A案", "B案"],
                    "rows": [["コスト", "低"]]}],
    })
    out = jp_slides.render_deck(content, out=str(tmp_path / "o.pptx"))
    prs = Presentation(out)
    t = [sh for sh in prs.slides[0].shapes if sh.has_table][0].table
    assert t.cell(1, 0).text == "コスト"
    assert t.cell(1, 1).text == "低"
    assert t.cell(1, 2).text == ""


def test_table_many_rows_height_capped(tmp_path):
    rows = [[f"項目{i}", f"値{i}"] for i in range(12)]
    content = jp_slides.parse_content({
        "theme": "藍",
        "slides": [{"type": "table", "headline": "長い表", "columns": ["項目", "値"], "rows": rows}],
    })
    out = jp_slides.render_deck(content, out=str(tmp_path / "o.pptx"))
    prs = Presentation(out)
    shp = [sh for sh in prs.slides[0].shapes if sh.has_table][0]
    assert shp.top + shp.height <= 6858000  # 7.5in in EMU


# ---- 画像: パス欠損は警告してスキップ ----

def test_image_missing_path_warns_but_builds(tmp_path, capsys):
    content = jp_slides.parse_content({
        "theme": "藍",
        "slides": [{"type": "image", "headline": "無い画像",
                    "image": str(tmp_path / "nope.png"), "caption": "図2"}],
    })
    out = jp_slides.render_deck(content, out=str(tmp_path / "o.pptx"))
    err = capsys.readouterr().err
    assert "画像" in err
    prs = Presentation(out)
    assert len(prs.slides) == 1


def test_image_valid_path_inserts_picture(tmp_path):
    png = tmp_path / "x.png"
    png.write_bytes(bytes.fromhex(_ONE_PX_PNG_HEX))
    content = jp_slides.parse_content({
        "theme": "藍",
        "slides": [{"type": "image", "headline": "推移", "image": str(png), "caption": "図1"}],
    })
    out = jp_slides.render_deck(content, out=str(tmp_path / "o.pptx"))
    prs = Presentation(out)
    pics = [sh for sh in prs.slides[0].shapes if sh.shape_type == 13]  # PICTURE
    assert pics, "画像が挿入されていない"


# ---- gallery / overview ----

def test_gallery_makes_16_plus_overview(tmp_path):
    paths = jp_slides.build_gallery(str(tmp_path / "gal"))
    decks = list((tmp_path / "gal").glob("*.pptx"))
    assert len([p for p in decks if p.name != "overview.pptx"]) == 16
    assert (tmp_path / "gal" / "overview.pptx").exists()
    assert len(paths) == 17


def test_overview_has_16_theme_names_text(tmp_path):
    out = jp_slides.build_overview(str(tmp_path / "ov.pptx"))
    prs = Presentation(out)
    text = _all_text(prs)
    for key in jp_slides.THEME_PROFILES:
        assert key in text


def test_gallery_cli(tmp_path):
    rc = jp_slides.main(["jp_slides.py", "gallery", "--out-dir", str(tmp_path / "g2")])
    assert rc == 0
    assert (tmp_path / "g2" / "overview.pptx").exists()


def test_overview_cli(tmp_path):
    outp = tmp_path / "o.pptx"
    rc = jp_slides.main(["jp_slides.py", "overview", "--out", str(outp)])
    assert rc == 0 and outp.exists()
