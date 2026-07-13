import pytest

pytest.importorskip("pptx")

import jp_slides
import jp_slides_design as D
from pptx import Presentation

THEME_KEYS = list(jp_slides.THEME_PROFILES)

# 各レイアウトの代表テーマ(タスク指示どおり)
LIST_KEY = "白磁"  # ※ render_stats の layout は poster だが、message の bullet 代表として使用
BENTO_KEY = "藍"
CARDS_KEY = "常磐"
POSTER_KEY = "朱"
LIST_LAYOUT_KEY = "霞"  # layout="list" の実テーマ


def _all_text(slide):
    chunks = []
    for shp in slide.shapes:
        if shp.has_text_frame:
            chunks.append(shp.text_frame.text)
    return "\n".join(chunks)


def _build(fn, key, *args):
    prof = jp_slides.THEME_PROFILES[key]
    prs = D.new_prs()
    slide = fn(prs, prof, *args)
    return prs, slide


# ---- render_message ----

@pytest.mark.parametrize("key", THEME_KEYS)
def test_render_message_builds_valid_pptx_all_themes(tmp_path, key):
    prs, slide = _build(D.render_message, key, "主要指標は計画を上回った",
                         ["新規契約が8件成立", "解約はゼロ", "主力商品Aが牽引"])
    out = tmp_path / f"{key}.pptx"
    prs.save(str(out))
    reopened = Presentation(str(out))
    assert len(reopened.slides) == 1


@pytest.mark.parametrize("key", ["霞", "藤", "明朝", "常磐", "藍", "墨"])
def test_render_message_headline_and_body_text_present(key):
    _, slide = _build(D.render_message, key, "見出しテキスト",
                       ["一つ目の要点", "二つ目の要点", "三つ目の要点"])
    text = _all_text(slide)
    assert "見出しテキスト" in text
    assert "一つ目の要点" in text
    assert "二つ目の要点" in text
    assert "三つ目の要点" in text


def test_render_message_empty_body_is_safe():
    _, slide = _build(D.render_message, "藍", "見出しのみ", [])
    assert "見出しのみ" in _all_text(slide)


def test_render_message_none_body_is_safe():
    _, slide = _build(D.render_message, "藍", "見出しのみ", None)
    assert "見出しのみ" in _all_text(slide)


def test_render_message_long_body_builds_without_crash(tmp_path):
    body = [f"論点その{i}" for i in range(1, 7)]
    prs, slide = _build(D.render_message, "霞", "六つの論点", body)
    out = tmp_path / "long.pptx"
    prs.save(str(out))
    text = _all_text(slide)
    for item in body:
        assert item in text


@pytest.mark.parametrize("key,bullet", [
    ("石板", "square"), ("白磁", "dash"), ("常磐", "chevron"),
    ("藍鉄", "ring"), ("山吹", "number"), ("明朝", "tick"),
])
def test_render_message_bullet_style_matches_profile(key, bullet):
    prof = jp_slides.THEME_PROFILES[key]
    assert prof["bullet"] == bullet
    # スタイル別にクラッシュしないことのみ確認(図形種別までは必須検証しない)
    _, slide = _build(D.render_message, key, "見出し", ["項目A", "項目B"])
    assert "項目A" in _all_text(slide)


# ---- render_stats ----

@pytest.mark.parametrize("key", THEME_KEYS)
def test_render_stats_smoke_all_themes(tmp_path, key):
    items = [
        {"value": "182億", "label": "エネルギー", "note": "+18%"},
        {"value": "24億", "label": "SaaS", "note": "黒字化"},
        {"value": "97億", "label": "海外", "note": "+4%"},
    ]
    prs, slide = _build(D.render_stats, key, "主要指標の伸び", items)
    out = tmp_path / f"{key}.pptx"
    prs.save(str(out))
    reopened = Presentation(str(out))
    assert len(reopened.slides) == 1


def test_render_stats_bento_layout_shows_all_values_and_labels():
    items = [
        {"value": "182億", "label": "エネルギー", "note": "+18%"},
        {"value": "24億", "label": "SaaS", "note": "黒字化"},
        {"value": "97億", "label": "海外", "note": "+4%"},
    ]
    _, slide = _build(D.render_stats, BENTO_KEY, "三つの事業", items)
    text = _all_text(slide)
    for it in items:
        assert it["value"] in text
        assert it["label"] in text


def test_render_stats_cards_layout_shows_all_values_and_labels():
    items = [
        {"value": "182億", "label": "エネルギー", "note": "+18%"},
        {"value": "24億", "label": "SaaS", "note": "黒字化"},
        {"value": "97億", "label": "海外", "note": "+4%"},
    ]
    _, slide = _build(D.render_stats, CARDS_KEY, "三つの事業", items)
    text = _all_text(slide)
    for it in items:
        assert it["value"] in text
        assert it["label"] in text


def test_render_stats_poster_single_item_oversized_value():
    items = [{"value": "×2.0", "label": "認知度", "note": "3か月で倍増"}]
    _, slide = _build(D.render_stats, POSTER_KEY, "ブランド刷新の成果", items)
    text = _all_text(slide)
    assert "×2.0" in text
    assert "認知度" in text
    assert "3か月で倍増" in text
    sizes = [r.font.size.pt for shp in slide.shapes if shp.has_text_frame
             for p in shp.text_frame.paragraphs for r in p.runs if r.text == "×2.0"]
    assert sizes and sizes[0] >= 130


def test_render_stats_poster_multi_item_lists_numbers():
    items = [
        {"value": "14%", "label": "売上増"},
        {"value": "800", "label": "新規顧客"},
    ]
    _, slide = _build(D.render_stats, POSTER_KEY, "数字で見る成果", items)
    text = _all_text(slide)
    assert "14%" in text
    assert "800" in text


def test_render_stats_list_layout_shows_all_rows():
    items = [
        {"value": "182億", "label": "エネルギー", "note": "+18%"},
        {"value": "24億", "label": "SaaS"},
    ]
    _, slide = _build(D.render_stats, LIST_LAYOUT_KEY, "主要指標", items)
    text = _all_text(slide)
    assert "182億" in text
    assert "エネルギー" in text
    assert "24億" in text
    assert "SaaS" in text


def test_render_stats_empty_items_is_safe():
    _, slide = _build(D.render_stats, "藍", "見出しのみ", [])
    assert "見出しのみ" in _all_text(slide)


def test_render_stats_none_items_is_safe():
    _, slide = _build(D.render_stats, "藍", "見出しのみ", None)
    assert "見出しのみ" in _all_text(slide)


def test_render_stats_long_items_builds_without_crash(tmp_path):
    items = [{"value": f"{i}0億", "label": f"事業{i}"} for i in range(1, 7)]
    prs, slide = _build(D.render_stats, CARDS_KEY, "六つの事業", items)
    out = tmp_path / "long_stats.pptx"
    prs.save(str(out))
    text = _all_text(slide)
    for it in items:
        assert it["value"] in text


@pytest.mark.parametrize("key", THEME_KEYS)
def test_render_stats_smoke_three_items_all_16_themes(key):
    items = [
        {"value": "12", "label": "指標A"},
        {"value": "34", "label": "指標B"},
        {"value": "56", "label": "指標C"},
    ]
    _, slide = _build(D.render_stats, key, "スモークテスト", items)
    text = _all_text(slide)
    assert "スモークテスト" in text
