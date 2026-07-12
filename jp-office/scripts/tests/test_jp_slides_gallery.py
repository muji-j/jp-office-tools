from pathlib import Path
import pytest
import jp_slides
pptx = pytest.importorskip("pptx")
from pptx import Presentation


def test_gallery_makes_16_plus_overview(tmp_path):
    paths = jp_slides.build_gallery(str(tmp_path / "gal"))
    decks = list((tmp_path / "gal").glob("*.pptx"))
    assert len([p for p in decks if p.name != "overview.pptx"]) == 16
    assert (tmp_path / "gal" / "overview.pptx").exists()


def test_overview_has_16_cards_text(tmp_path):
    out = jp_slides.build_overview(str(tmp_path / "ov.pptx"))
    prs = Presentation(out)
    text = "\n".join(sh.text_frame.text for s in prs.slides for sh in s.shapes if sh.has_text_frame)
    for key in jp_slides.THEMES:
        assert key in text


def test_gallery_cli(tmp_path):
    rc = jp_slides.main(["jp_slides.py", "gallery", "--out-dir", str(tmp_path / "g2")])
    assert rc == 0
    assert (tmp_path / "g2" / "overview.pptx").exists()


def test_overview_cli(tmp_path):
    outp = tmp_path / "o.pptx"
    rc = jp_slides.main(["jp_slides.py", "overview", "--out", str(outp)])
    assert rc == 0 and outp.exists()
